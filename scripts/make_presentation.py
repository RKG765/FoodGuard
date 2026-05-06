"""
FoodGuard — Professor-Ready Technical Presentation
====================================================
Includes real eval results, deep ML justification slides,
and "why X not Y" reasoning for every design choice.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x0D, 0x0D, 0x0D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x00, 0xC8, 0x53)
RED        = RGBColor(0xFF, 0x17, 0x44)
AMBER      = RGBColor(0xFF, 0xAB, 0x00)
BLUE       = RGBColor(0x00, 0x96, 0xFF)
GREY_L     = RGBColor(0xF5, 0xF5, 0xF5)
GREY_M     = RGBColor(0xBD, 0xBD, 0xBD)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
NAVY       = RGBColor(0x0F, 0x34, 0x60)
SW, SH     = Inches(13.33), Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def _box(s, l, t, w, h, fc=None, bc=None, bp=0):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    if fc: sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:  sh.fill.background()
    if bc and bp: sh.line.color.rgb = bc; sh.line.width = Pt(bp)
    else: sh.line.fill.background()
    return sh

def _t(s, text, l, t, w, h, sz=28, b=False, c=BLACK, a=PP_ALIGN.LEFT, wrap=True, i=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = a
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = b; r.font.italic = i
    r.font.color.rgb = c; r.font.name = "Segoe UI"
    return tb

def _line(s, l, t, w, c=GREEN, th=4):
    ln = s.shapes.add_shape(1, l, t, w, Pt(th))
    ln.fill.solid(); ln.fill.fore_color.rgb = c; ln.line.fill.background()

def _header(s, label, color, title, dark=False):
    _line(s, Inches(0.6), Inches(0.55), Inches(1.2), color)
    _t(s, label, Inches(0.6), Inches(0.65), Inches(8), Inches(0.5), sz=13, b=True, c=color)
    _t(s, title, Inches(0.6), Inches(1.05), Inches(11.5), Inches(0.7),
       sz=36, b=True, c=WHITE if dark else BLACK)

def _stat(s, l, t, w, h, num, lbl, nc=GREEN, bg_=WHITE):
    _box(s, l, t, w, h, fc=bg_, bc=GREY_M, bp=1)
    _t(s, num, l, t+Inches(0.12), w, Inches(0.7), sz=48, b=True, c=nc, a=PP_ALIGN.CENTER)
    _t(s, lbl, l, t+Inches(0.82), w, Inches(0.45), sz=13, c=RGBColor(0x55,0x55,0x55), a=PP_ALIGN.CENTER)

# ======================= SLIDES ==========================

def s01_title(prs):
    s = _blank(prs); _bg(s, BLACK)
    _line(s, Inches(0), Inches(0), SW, GREEN, 6)
    _t(s, "🛡️", Inches(0), Inches(0.6), SW, Inches(1.2), sz=72, a=PP_ALIGN.CENTER)
    _t(s, "FoodGuard", Inches(0), Inches(1.6), SW, Inches(1.5), sz=80, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    _t(s, "AI-Powered Food Image Fraud Detection System", Inches(0), Inches(3.0), SW, Inches(0.8), sz=28, c=GREEN, a=PP_ALIGN.CENTER)
    _t(s, "Deep Learning Forensics  ·  EfficientNet-B3  ·  Stable Diffusion XL  ·  Threshold-Calibrated Inference",
       Inches(0), Inches(3.85), SW, Inches(0.6), sz=15, c=GREY_M, a=PP_ALIGN.CENTER)
    _t(s, "Raj  ·  Rahul  ·  Aman\nBML  |  Semester VI  |  Project 3  |  March 2026",
       Inches(0), Inches(6.5), SW, Inches(0.8), sz=13, c=GREY_M, a=PP_ALIGN.CENTER)

def s02_problem(prs):
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "PROBLEM STATEMENT", RED, "Visual Fraud in the Digital Food Ecosystem")
    points = [
        ("Fake Complaint Fraud", "Users digitally insert cockroaches, hair, or mold into real food photos to claim refunds from delivery platforms (Zomato, Swiggy, DoorDash)."),
        ("AI-Generated Deception", "Text-to-image models like Stable Diffusion can now produce photorealistic food images indistinguishable from real photographs to the human eye."),
        ("Binary Classifier Failure", "Traditional real-vs-fake classifiers fail when 98% of the image is genuine and only a 2-4% inpainted patch is fraudulent — the global image statistics remain 'real'."),
    ]
    for idx, (title, desc) in enumerate(points):
        y = Inches(2.2 + idx * 1.4)
        _box(s, Inches(0.5), y, Inches(12.0), Inches(1.2), fc=GREY_L, bc=RED, bp=2)
        _t(s, title, Inches(0.8), y + Inches(0.1), Inches(3.0), Inches(0.5), sz=17, b=True, c=RED)
        _t(s, desc, Inches(0.8), y + Inches(0.55), Inches(11.5), Inches(0.6), sz=14, c=RGBColor(0x33,0x33,0x33))
    _t(s, "Objective: Build a 4-class detector with FPR ≤ 5% on real images — genuine photos must never be wrongly flagged.",
       Inches(0.6), Inches(6.5), Inches(11), Inches(0.6), sz=15, b=True, c=BLUE, i=True)

def s03_why_4class(prs):
    """Why 4 classes, not binary — deep justification."""
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "DESIGN DECISION", AMBER, "Why 4 Classes Instead of Binary?")
    # Binary side
    _box(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5), fc=GREY_L, bc=RED, bp=3)
    _t(s, "❌  Binary Classifier (Real vs Fake)", Inches(0.7), Inches(2.1), Inches(5.4), Inches(0.6), sz=18, b=True, c=RED)
    _t(s, "• Loses the WHY — cannot distinguish a fully AI image from a subtly inpainted one\n"
          "• Treats a 100% synthetic image identically to a 2% tampered one\n"
          "• Cannot inform investigators what type of manipulation occurred\n"
          "• Global feature statistics of an inpainted image look 'real' — misses local anomalies",
       Inches(0.7), Inches(2.8), Inches(5.4), Inches(3.5), sz=14, c=RGBColor(0x33,0x33,0x33))
    # 4-class side
    _box(s, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5), fc=GREY_L, bc=GREEN, bp=3)
    _t(s, "✅  4-Class Classifier (Our Approach)", Inches(7.0), Inches(2.1), Inches(5.4), Inches(0.6), sz=18, b=True, c=GREEN)
    _t(s, "• Detects fully AI-generated images (texture-level artifacts)\n"
          "• Handles compression-degraded AI (social media sharing scenario)\n"
          "• Catches subtle inpainting (2-4% altered area — deliberate fraud)\n"
          "• Actionable output: tells you the specific type of manipulation\n"
          "• Can be extended to more classes via transfer learning",
       Inches(7.0), Inches(2.8), Inches(5.4), Inches(3.5), sz=14, c=RGBColor(0x33,0x33,0x33))

def s04_data_sources(prs):
    s = _blank(prs); _bg(s, DARK)
    _line(s, Inches(0), Inches(0), SW, BLUE, 6)
    _header(s, "DATA COLLECTION", BLUE, "191,000+ Real Images from 3 Kaggle Datasets", dark=True)
    datasets = [
        ("Food-101 (ETH Zurich)", "101,000 images", "101 categories — Western & international\nHigh-quality studio-like images\nMost widely used food benchmark", GREEN),
        ("Indian Food Dataset", "4,000+ images", "Indian cuisines (biryani, paneer, samosa)\nFills Food-101's blind spot on Indian food\nTarget demographic: Indian delivery apps", AMBER),
        ("UECFOOD256 + AIcrowd", "86,000 images", "Japanese food (256 categories)\nReal-world noisy photos (phone cameras)\nIncludes blurry, poorly lit images — critical for robustness", BLUE),
    ]
    for i, (name, count, desc, col) in enumerate(datasets):
        x = Inches(0.5 + i * 4.2)
        _box(s, x, Inches(2.3), Inches(3.9), Inches(4.5), fc=NAVY)
        _box(s, x, Inches(2.3), Inches(3.9), Inches(0.1), fc=col)
        _t(s, count, x, Inches(2.5), Inches(3.9), Inches(0.7), sz=32, b=True, c=GREEN, a=PP_ALIGN.CENTER)
        _t(s, name, x, Inches(3.2), Inches(3.9), Inches(0.5), sz=14, b=True, c=WHITE, a=PP_ALIGN.CENTER)
        _t(s, desc, x + Inches(0.2), Inches(3.8), Inches(3.5), Inches(2.8), sz=13, c=GREY_M)

def s05_ai_generation(prs):
    """Technical details of SDXL generation pipeline."""
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "ADVERSARIAL DATA SYNTHESIS", BLUE, "AI Image Generation Pipeline")
    # Left: Text-to-Image
    _box(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0), fc=GREY_L, bc=BLUE, bp=2)
    _t(s, "Text-to-Image (Classes 1, 2, 3)", Inches(0.7), Inches(2.1), Inches(5.6), Inches(0.5), sz=17, b=True, c=BLUE)
    _t(s, "Model: RealVisXL V4.0 (SG161222/RealVisXL_V4.0)\n"
          "Architecture: Stable Diffusion XL fine-tune (~3.5B params)\n\n"
          "Generation Parameters:\n"
          "  • Inference Steps: 25 (DDIM Scheduler)\n"
          "  • Guidance Scale: 5.0 – 7.5\n"
          "  • Resolution: 512 × 512\n"
          "  • Prompt engineering: 12 cuisines × quality modifiers\n\n"
          "Post-Processing for Class 2 (Compressed):\n"
          "  • Random downscale to 384/448/640px → rescale to 512\n"
          "  • JPEG compression: quality 40–85\n\n"
          "Post-Processing for Class 3 (Degraded):\n"
          "  • Gaussian blur (r=0.3–0.8, 70% probability)\n"
          "  • Gaussian noise (σ = 1.5–5.0)",
       Inches(0.7), Inches(2.7), Inches(5.6), Inches(4.0), sz=12, c=RGBColor(0x33,0x33,0x33))
    # Right: Inpainting
    _box(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(5.0), fc=GREY_L, bc=RED, bp=2)
    _t(s, "Inpainting Pipeline (Class 4 — Fraud)", Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.5), sz=17, b=True, c=RED)
    _t(s, "Model: RealVisXL V4.0 via StableDiffusionXLInpaintPipeline\n\n"
          "Mask Generation Algorithm:\n"
          "  • 2–3 overlapping ellipses (irregular blob)\n"
          "  • Coverage: 2–4% of image area only\n"
          "  • Center-biased placement (30–70% of dims)\n"
          "  • Gaussian blur on edges for seamless blending\n\n"
          "Inpainting Parameters:\n"
          "  • CFG Scale: 4.5 (very low — prevents object dominance)\n"
          "  • Denoising Strength: 0.99\n"
          "  • Steps: 26\n\n"
          "Fraud Objects (weighted distribution):\n"
          "  cockroach, housefly, mosquito, bee, ant, worm,\n"
          "  human hair, mold, plastic, paper, metal shard",
       Inches(7.0), Inches(2.7), Inches(5.6), Inches(4.0), sz=12, c=RGBColor(0x33,0x33,0x33))

def s06_why_realvis(prs):
    """Why RealVisXL and not DALL-E / Flux / SD1.5"""
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "DESIGN DECISION", AMBER, "Why RealVisXL V4.0 and Not Other Models?")
    rows = [
        ("RealVisXL V4.0 [OK]", "Open-source, free, runs locally on 12GB VRAM.\nFine-tuned specifically for photorealism — critical for training a robust detector.", GREEN),
        ("DALL-E 3 [FAIL]", "Closed-source, API-only ($0.04/image). 2000 images = $80+.\nNot scriptable for batch generation. Rate limited.", RED),
        ("Midjourney [FAIL]", "Closed-source, Discord-only interface.\nCannot automate via Python scripts.", RED),
        ("Stable Diffusion 1.5 [FAIL]", "512×512 but significantly lower quality.\nAI artifacts too obvious — detector would have an unfair advantage (not realistic evaluation).", RED),
        ("Flux / SD3 [FAIL]", "Requires 16+ GB VRAM — exceeds our RTX 5070 Ti (12GB).\nLess tested with SDXL Inpainting Pipeline.", RED),
    ]
    for i, (model, reason, col) in enumerate(rows):
        y = Inches(2.0 + i * 1.05)
        _box(s, Inches(0.5), y, Inches(12.0), Inches(0.95), fc=GREY_L, bc=col, bp=2)
        _box(s, Inches(0.5), y, Inches(0.1), Inches(0.95), fc=col)
        _t(s, model, Inches(0.8), y + Inches(0.1), Inches(3.0), Inches(0.8), sz=15, b=True, c=col)
        _t(s, reason, Inches(3.8), y + Inches(0.05), Inches(8.5), Inches(0.9), sz=12, c=RGBColor(0x33,0x33,0x33))

def s07_model_arch(prs):
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "MODEL ARCHITECTURE", GREEN, "EfficientNet-B3 — Forward Pass Pipeline")
    # Flow boxes
    flow = [
        ("Input Image\n512 × 512 × 3", GREY_L),
        ("EfficientNet-B3\nPretrained ImageNet\n~12M params", NAVY),
        ("Global Avg Pool\n1536-dim vector", GREY_L),
        ("Linear Head\n1536 → 4", GREY_L),
        ("Softmax\n4 probabilities", GREY_L),
    ]
    for i, (label, fc) in enumerate(flow):
        x = Inches(0.5 + i * 2.45)
        tc = WHITE if fc == NAVY else BLACK
        _box(s, x, Inches(2.3), Inches(2.2), Inches(1.4), fc=fc, bc=GREEN, bp=2)
        _t(s, label, x, Inches(2.4), Inches(2.2), Inches(1.2), sz=13, b=True, a=PP_ALIGN.CENTER, c=tc)
        if i < len(flow) - 1:
            _t(s, "→", x + Inches(2.2), Inches(2.7), Inches(0.3), Inches(0.5), sz=28, c=GREEN)
    # Decision logic
    _box(s, Inches(0.5), Inches(4.2), Inches(12.0), Inches(1.4), fc=GREY_L, bc=AMBER, bp=2)
    _t(s, "Threshold Decision Logic", Inches(0.7), Inches(4.3), Inches(5), Inches(0.4), sz=16, b=True, c=AMBER)
    _t(s, "IF  P(real) > θ   →   Classify as REAL\n"
          "ELSE              →   Classify as argmax(P(perfect_ai), P(compressed_ai), P(edited_ai))\n"
          "θ calibrated on validation set to guarantee FPR ≤ 5%",
       Inches(0.7), Inches(4.75), Inches(11.5), Inches(0.8), sz=14, c=RGBColor(0x33,0x33,0x33))

    _t(s, "AMP Optimization: torch.cuda.amp enables FP16/FP32 mixed precision — 40% VRAM reduction, allowing batch_size=16 for 512×512 on 12GB GPU.",
       Inches(0.5), Inches(6.3), Inches(12), Inches(0.5), sz=13, i=True, c=GREY_M)

def s08_why_effnet(prs):
    """Why EfficientNet-B3 and not ResNet/ViT/B0/B4"""
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "DESIGN DECISION", AMBER, "Why EfficientNet-B3 and Not Other Backbones?")
    rows = [
        ("EfficientNet-B3 [OK]", "Best accuracy-per-parameter via neural architecture search (compound scaling).\n12M params, fits 512×512 @ batch=16 in 12GB VRAM.", GREEN),
        ("EfficientNet-B0/B1 [FAIL]", "Too few parameters (~5M) for 512×512 forensic feature capture.\nWould miss subtle high-frequency AI artifacts in textures.", RED),
        ("EfficientNet-B4+ [FAIL]", "Exceeds VRAM budget: B4 @ 512×512 requires batch_size < 8.\nSlows training significantly without proportional accuracy gain.", RED),
        ("ResNet-50 [FAIL]", "25M params but similar accuracy to B3. Less parameter-efficient.\nNo compound scaling — manual depth/width trade-offs.", RED),
        ("Vision Transformer (ViT) [FAIL]", "Requires 100K+ labeled images per class to outperform CNNs.\nWith only ~7K total images, ViT would severely overfit.", RED),
    ]
    for i, (model, reason, col) in enumerate(rows):
        y = Inches(2.0 + i * 1.05)
        _box(s, Inches(0.5), y, Inches(12.0), Inches(0.95), fc=GREY_L, bc=col, bp=2)
        _box(s, Inches(0.5), y, Inches(0.1), Inches(0.95), fc=col)
        _t(s, model, Inches(0.8), y + Inches(0.1), Inches(3.2), Inches(0.8), sz=15, b=True, c=col)
        _t(s, reason, Inches(4.0), y + Inches(0.05), Inches(8.3), Inches(0.9), sz=12, c=RGBColor(0x33,0x33,0x33))

def s09_training_config(prs):
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "TRAINING CONFIGURATION", BLUE, "Hyperparameters & Optimization Strategy")
    # Left: Optimizer & Loss
    _box(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0), fc=GREY_L, bc=AMBER, bp=2)
    _t(s, "Loss Function & Optimizer", Inches(0.7), Inches(2.1), Inches(5.6), Inches(0.5), sz=17, b=True, c=AMBER)
    _t(s, "Loss: Weighted Cross-Entropy\n"
          "  Weights: [1.2, 1.0, 1.0, 1.0]\n"
          "  Real class penalized 20% more → reduces FPR\n\n"
          "Why not Focal Loss?\n"
          "  → Focal Loss implemented as backup (focal_loss.py)\n"
          "  → CE converged well; Focal can destabilize if γ not tuned\n\n"
          "Optimizer: AdamW\n"
          "  lr = 3e-4, weight_decay = 1e-4\n\n"
          "Why AdamW, not SGD?\n"
          "  → AdamW decouples weight decay from gradients\n"
          "  → More forgiving for fine-tuning pretrained models\n"
          "  → SGD needs careful warmup and manual LR scheduling",
       Inches(0.7), Inches(2.7), Inches(5.6), Inches(4.0), sz=12, c=RGBColor(0x33,0x33,0x33))
    # Right: Scheduler & Data
    _box(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(5.0), fc=GREY_L, bc=BLUE, bp=2)
    _t(s, "Scheduler, Data & Augmentation", Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.5), sz=17, b=True, c=BLUE)
    _t(s, "Scheduler: Cosine Annealing (T_max = 20)\n"
          "  Smooth LR decay, no sudden drops\n\n"
          "Why not StepLR or ReduceOnPlateau?\n"
          "  → StepLR drops LR abruptly — destabilizes training\n"
          "  → ReduceOnPlateau is reactive; Cosine is proactive\n\n"
          "Data Split: 70% Train / 15% Val / 15% Test\n"
          "  Train: 4,947  |  Val: 1,060  |  Test: 1,061\n\n"
          "Augmentation: MINIMAL by design\n"
          "  Only: RandomHorizontalFlip + ColorJitter(0.1)\n\n"
          "Why minimal?\n"
          "  → Heavy blur/noise would DESTROY the AI artifacts\n"
          "  → Class 2 & 3 already USE compression as features\n"
          "  → Adding random blur confuses the detector",
       Inches(7.0), Inches(2.7), Inches(5.6), Inches(4.0), sz=12, c=RGBColor(0x33,0x33,0x33))

def s10_why_512(prs):
    """Why 512x512 resolution"""
    s = _blank(prs); _bg(s, DARK)
    _line(s, Inches(0), Inches(0), SW, GREEN, 6)
    _header(s, "DESIGN DECISION", GREEN, "Why 512×512 Resolution?", dark=True)
    rows = [
        ("224 × 224", "Standard ImageNet size. Loses too much forensic detail.\nAI artifacts exist at pixel level (texture, edge smoothing) — downsampling destroys them.", RED),
        ("512 × 512 [OK]", "Sweet spot: preserves high-frequency forensic signals while fitting in 12GB VRAM.\nBatch=16 possible with AMP. Matches our SDXL generation resolution.", GREEN),
        ("1024 × 1024", "Ideal for forensics but requires >24GB VRAM with EfficientNet-B3 @ batch=16.\nNot feasible on RTX 5070 Ti without extreme batch reduction.", RED),
    ]
    for i, (res, reason, col) in enumerate(rows):
        y = Inches(2.2 + i * 1.5)
        _box(s, Inches(0.5), y, Inches(12.0), Inches(1.3), fc=NAVY)
        _box(s, Inches(0.5), y, Inches(0.12), Inches(1.3), fc=col)
        _t(s, res, Inches(0.8), y + Inches(0.15), Inches(2.5), Inches(0.5), sz=22, b=True, c=col)
        _t(s, reason, Inches(3.5), y + Inches(0.1), Inches(9.0), Inches(1.1), sz=14, c=GREY_M)

def s11_threshold(prs):
    s = _blank(prs); _bg(s, DARK)
    _line(s, Inches(0), Inches(0), SW, GREEN, 6)
    _header(s, "EVALUATION METHODOLOGY", GREEN, "Threshold Calibration for Strict FPR Control", dark=True)
    _t(s, "Enterprise Requirement: False Positive Rate ≤ 5%\nReal food photos must NEVER be wrongly flagged as AI.",
       Inches(0.6), Inches(2.1), Inches(12), Inches(0.9), sz=18, c=GREY_M)
    _box(s, Inches(0.5), Inches(3.2), Inches(12.0), Inches(2.0), fc=NAVY)
    _t(s, "Algorithm:\n"
          "1. After training, compute Softmax: [P(real), P(perfect), P(compressed), P(edited)]\n"
          "2. Sweep threshold θ from 0.50 to 0.99 (100 steps) on validation set\n"
          "3. Decision rule: IF P(real) > θ → REAL,  ELSE → argmax of AI classes\n"
          "4. Select θ that satisfies FPR ≤ 0.05 with maximum overall accuracy",
       Inches(0.7), Inches(3.4), Inches(11.6), Inches(1.6), sz=15, c=WHITE)
    _t(s, "Result: Optimal threshold = 0.50 (default softmax boundary).\n"
          "The network separated feature spaces so cleanly that no forced thresholding was needed.",
       Inches(0.6), Inches(5.6), Inches(12), Inches(1.0), sz=16, i=True, c=GREEN)

def s12_results_hero(prs):
    """Big 0% FPR slide"""
    s = _blank(prs); _bg(s, BLACK)
    _line(s, Inches(0), Inches(0), SW, GREEN, 6)
    _t(s, "TEST SET RESULTS", Inches(0), Inches(0.6), SW, Inches(0.5), sz=14, b=True, c=GREEN, a=PP_ALIGN.CENTER)
    _t(s, "0.00%", Inches(0), Inches(1.2), SW, Inches(2.5), sz=160, b=True, c=GREEN, a=PP_ALIGN.CENTER)
    _t(s, "False Positive Rate on Real Images", Inches(0), Inches(3.6), SW, Inches(0.6), sz=26, c=WHITE, a=PP_ALIGN.CENTER)
    _t(s, "Not a single genuine food photo was wrongly flagged as AI.", Inches(0), Inches(4.3), SW, Inches(0.6), sz=18, i=True, c=GREY_M, a=PP_ALIGN.CENTER)
    _stat(s, Inches(1.5), Inches(5.2), Inches(2.5), Inches(1.5), "99.81%", "Test Accuracy", GREEN, DARK)
    _stat(s, Inches(4.3), Inches(5.2), Inches(2.5), Inches(1.5), "0.0004", "Final Loss", BLUE, DARK)
    _stat(s, Inches(7.1), Inches(5.2), Inches(2.5), Inches(1.5), "Epoch 18", "Best Model", AMBER, DARK)
    _stat(s, Inches(9.9), Inches(5.2), Inches(2.5), Inches(1.5), "θ = 0.50", "Threshold", GREEN, DARK)
    for sh in s.shapes:
        for p in (sh.text_frame.paragraphs if sh.has_text_frame else []):
            for r in p.runs:
                if r.font.color.rgb == RGBColor(0x55,0x55,0x55): r.font.color.rgb = GREY_M

def s13_eval_report(prs):
    """Full classification report from evaluate.py"""
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "EVALUATION REPORT", BLUE, "Classification Report (evaluate.py Output)")
    # Table
    headers = ["CLASS", "PRECISION", "RECALL", "F1-SCORE", "SUPPORT"]
    cx = [Inches(0.5), Inches(3.5), Inches(5.5), Inches(7.5), Inches(9.5)]
    cw = [Inches(2.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.0)]
    _box(s, Inches(0.4), Inches(2.0), Inches(12.0), Inches(0.5), fc=DARK)
    for j, h in enumerate(headers):
        _t(s, h, cx[j], Inches(2.05), cw[j], Inches(0.4), sz=12, b=True, c=WHITE)
    rows = [
        ("Real",          "1.0000", "1.0000", "1.0000",  "90", GREEN),
        ("Perfect AI",    "0.9917", "0.9917", "0.9917", "120", BLUE),
        ("Compressed AI", "1.0000", "1.0000", "1.0000", "101", AMBER),
        ("Edited AI",     "0.9987", "0.9987", "0.9987", "750", RED),
    ]
    for i, (cls, prec, rec, f1, sup, col) in enumerate(rows):
        y = Inches(2.55 + i * 0.75)
        bg_ = GREY_L if i % 2 == 0 else WHITE
        _box(s, Inches(0.4), y, Inches(12.0), Inches(0.7), fc=bg_)
        _box(s, Inches(0.4), y, Inches(0.1), Inches(0.7), fc=col)
        vals = [cls, prec, rec, f1, sup]
        for j, v in enumerate(vals):
            _t(s, v, cx[j], y + Inches(0.1), cw[j], Inches(0.5), sz=14, b=(j==0), c=BLACK)
    # Aggregates
    y_agg = Inches(5.65)
    _box(s, Inches(0.4), y_agg, Inches(12.0), Inches(0.5), fc=DARK)
    aggs = [("Weighted Avg", "0.9981", "0.9981", "0.9981", "1061")]
    for cls, prec, rec, f1, sup in aggs:
        vals = [cls, prec, rec, f1, sup]
        for j, v in enumerate(vals):
            _t(s, v, cx[j], y_agg + Inches(0.05), cw[j], Inches(0.4), sz=12, b=True, c=WHITE)

    _t(s, "Key Insight: The model achieved perfect precision and recall on Real and Compressed AI classes.\n"
          "Only 2 misclassifications total across 1,061 test images (both between perfect_ai ↔ edited_ai).",
       Inches(0.5), Inches(6.4), Inches(12), Inches(0.8), sz=14, i=True, c=GREY_M)

def s14_confusion(prs):
    """Confusion matrix in table form"""
    s = _blank(prs); _bg(s, WHITE)
    _header(s, "CONFUSION MATRIX", AMBER, "Test Set — 1,061 Images")
    # Predicted headers
    _t(s, "PREDICTED →", Inches(3.0), Inches(2.0), Inches(2), Inches(0.4), sz=11, b=True, c=GREY_M)
    ph = ["real", "perfect_ai", "compress_ai", "edited_ai"]
    for j, h in enumerate(ph):
        _t(s, h, Inches(3.5 + j * 2.2), Inches(2.4), Inches(2.0), Inches(0.4), sz=12, b=True, c=BLACK, a=PP_ALIGN.CENTER)
    _t(s, "TRUE ↓", Inches(0.5), Inches(2.9), Inches(1.5), Inches(0.4), sz=11, b=True, c=GREY_M)
    # Matrix rows
    matrix = [
        ("real",          GREEN, [("90", GREEN), ("0", BLACK), ("0", BLACK), ("0", BLACK)]),
        ("perfect_ai",    BLUE,  [("0", BLACK), ("119", GREEN), ("0", BLACK), ("1", RED)]),
        ("compressed_ai", AMBER, [("0", BLACK), ("0", BLACK), ("101", GREEN), ("0", BLACK)]),
        ("edited_ai",     RED,   [("0", BLACK), ("1", RED), ("0", BLACK), ("749", GREEN)]),
    ]
    for i, (label, col, cells) in enumerate(matrix):
        y = Inches(3.0 + i * 0.8)
        bg_ = GREY_L if i % 2 == 0 else WHITE
        _box(s, Inches(0.4), y, Inches(12.0), Inches(0.75), fc=bg_)
        _box(s, Inches(0.4), y, Inches(0.1), Inches(0.75), fc=col)
        _t(s, label, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.55), sz=14, b=True, c=col)
        for j, (val, vc) in enumerate(cells):
            _t(s, val, Inches(3.5 + j * 2.2), y + Inches(0.1), Inches(2.0), Inches(0.55),
               sz=20, b=True, c=vc, a=PP_ALIGN.CENTER)
    _t(s, "Only 2 errors in 1,061 test images — both between semantically similar AI sub-classes.\n"
          "Zero errors on Real class → FPR = 0.00%",
       Inches(0.5), Inches(6.4), Inches(12), Inches(0.8), sz=15, b=True, c=GREEN, i=True)

def s15_future(prs):
    s = _blank(prs); _bg(s, DARK)
    _line(s, Inches(0), Inches(0), SW, BLUE, 6)
    _t(s, "FUTURE WORK & ENHANCEMENTS", Inches(0), Inches(0.8), SW, Inches(0.6), sz=14, b=True, c=BLUE, a=PP_ALIGN.CENTER)
    nexts = [
        (BLUE, "Grad-CAM Explainability", "Gradient-weighted Class Activation Maps to visualize WHERE the model detects AI artifacts.\nCritical for forensic evidence — proves WHY an image is flagged."),
        (GREEN, "REST API (FastAPI)", "Wrap inference.py in an async REST endpoint.\nEnable real-time detection for food delivery platforms via simple HTTP POST."),
        (AMBER, "Dual-Stream Model (RGB + FFT)", "Already coded: dual_stream_detector.py + ela.py.\nCombine spatial (RGB) features with frequency domain (FFT) forensic features."),
        (RED, "Scale to Full Dataset", "Current: 7K prototype. Full corpus: 191K+ images.\nTest generalization with 50× more data and adversarial augmentation."),
    ]
    for i, (col, title, desc) in enumerate(nexts):
        x = Inches(0.5 + (i % 2) * 6.3)
        y = Inches(1.8 + (i // 2) * 2.5)
        _box(s, x, y, Inches(5.9), Inches(2.2), fc=NAVY)
        _box(s, x, y, Inches(0.12), Inches(2.2), fc=col)
        _t(s, title, x + Inches(0.3), y + Inches(0.15), Inches(5.4), Inches(0.5), sz=19, b=True, c=col)
        _t(s, desc, x + Inches(0.3), y + Inches(0.65), Inches(5.4), Inches(1.4), sz=13, c=GREY_M)

def s16_end(prs):
    s = _blank(prs); _bg(s, BLACK)
    _line(s, Inches(0), Inches(0), SW, GREEN, 6)
    _t(s, "🛡️ FoodGuard", Inches(0), Inches(2.2), SW, Inches(1.2), sz=64, b=True, c=WHITE, a=PP_ALIGN.CENTER)
    _t(s, "99.81% Accuracy  ·  0% False Positives  ·  10.7M Parameters",
       Inches(0), Inches(3.5), SW, Inches(0.6), sz=22, c=GREEN, a=PP_ALIGN.CENTER)
    _t(s, "Thank You  ·  Questions?", Inches(0), Inches(5.0), SW, Inches(0.6), sz=24, c=GREY_M, a=PP_ALIGN.CENTER)
    _t(s, "Raj  ·  Rahul  ·  Aman  |  BML Semester VI  |  March 2026",
       Inches(0), Inches(6.8), SW, Inches(0.5), sz=13, c=GREY_M, a=PP_ALIGN.CENTER)

# ── Build ────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    s01_title(prs)             # 1  – Hero
    s02_problem(prs)           # 2  – Problem & motivation
    s03_why_4class(prs)        # 3  – Why 4-class not binary
    s04_data_sources(prs)      # 4  – 3 Kaggle datasets
    s05_ai_generation(prs)     # 5  – SDXL generation pipeline (technical)
    s06_why_realvis(prs)       # 6  – Why RealVisXL not DALL-E/Flux
    s07_model_arch(prs)        # 7  – EfficientNet-B3 pipeline
    s08_why_effnet(prs)        # 8  – Why B3 not ResNet/ViT
    s09_training_config(prs)   # 9  – Hyperparameters + justifications
    s10_why_512(prs)           # 10 – Why 512×512
    s11_threshold(prs)         # 11 – Threshold calibration math
    s12_results_hero(prs)      # 12 – Giant 0% FPR
    s13_eval_report(prs)       # 13 – Full classification report
    s14_confusion(prs)         # 14 – Confusion matrix
    s15_future(prs)            # 15 – Future work
    s16_end(prs)               # 16 – Thank you

    out = "FoodGuard_Final_Presentation.pptx"
    prs.save(out)
    print(f"\n[OK] Saved: {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
