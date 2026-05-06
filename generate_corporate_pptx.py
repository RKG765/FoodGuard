import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation with a widescreen format (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Corporate Color Palette ---
COLOR_BG = RGBColor(250, 252, 255)         # Very light blue/gray tint
COLOR_PRIMARY = RGBColor(15, 32, 67)       # Deep Corporate Navy
COLOR_ACCENT = RGBColor(0, 102, 204)       # Bright Azure Blue
COLOR_TEXT = RGBColor(60, 64, 67)          # Dark Charcoal for readability
COLOR_TEXT_LIGHT = RGBColor(128, 134, 139) # Light Charcoal for subtext
COLOR_WHITE = RGBColor(255, 255, 255)      # White

def apply_background(slide):
    # Set solid background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header_footer(slide, title_text, include_divider=True):
    # Header Background Shape
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.fill.background()

    # Title Text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = 'Arial'

    # Accent Divider Line
    if include_divider:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(13.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()
        
    # Footer Shape
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.5))
    tf_f = footer.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "FoodGuard | AI Forensic System | Mid-Semester Review"
    p_f.font.size = Pt(12)
    p_f.font.color.rgb = COLOR_TEXT_LIGHT
    p_f.font.name = 'Arial'

def add_body_content(slide, bullet_points):
    # Precise positioning for text to avoid messy lining
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for idx, text in enumerate(bullet_points):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = text
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXT
        p.font.name = 'Arial'
        
        # Proper spacing & lining
        p.space_after = Pt(16)
        
        # Sub-bullet handler (simple indent if starts with '-')
        if text.strip().startswith("-"):
            p.level = 1
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXT_LIGHT

# ----------------- SLIDE 1: TITLE -----------------
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
apply_background(slide)

# Big centered shape
center_rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5)
)
center_rect.fill.solid()
center_rect.fill.fore_color.rgb = COLOR_PRIMARY
center_rect.line.fill.background()

accent_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.0), Inches(13.333), Inches(0.1)
)
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = COLOR_ACCENT
accent_line.line.fill.background()

# Title
tx = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.5))
tf = tx.text_frame
p = tf.paragraphs[0]
p.text = "FoodGuard"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
p2 = tf.add_paragraph()
p2.text = "AI-Powered Food Fraud Detection System"
p2.font.size = Pt(28)
p2.font.color.rgb = COLOR_ACCENT
p2.alignment = PP_ALIGN.CENTER

# Team members at bottom right
tx_team = slide.shapes.add_textbox(Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.5))
tf_t = tx_team.text_frame
p_t = tf_t.paragraphs[0]
p_t.text = "Project Team:\nRaj Kumar\nAman Yadav\nRahul"
p_t.font.size = Pt(20)
p_t.font.color.rgb = COLOR_TEXT
p_t.font.bold = True
p_t.alignment = PP_ALIGN.RIGHT

# ----------------- SLIDE 2: PROBLEM STATEMENT -----------------
s2 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s2)
add_header_footer(s2, "Problem Statement")

add_body_content(s2, [
    "The rise of ultra-realistic AI generation (e.g., SDXL) threatens visual verification systems.",
    "Critical Fraud Vectors in the Food Industry:",
    "- Customers generating fake contaminants (insects, mold) for refunds.",
    "- Competitors staging fake health code violations.",
    "- Businesses generating synthetic menu items lacking authenticity.",
    "Why Current Solutions Fail:",
    "- Binary detectors (Real vs. Fake) do not explain the nature of the forgery.",
    "- They struggle with subtle, localized tampering (inpainting)."
])

# ----------------- SLIDE 3: OUR SOLUTION -----------------
s3 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s3)
add_header_footer(s3, "Our Solution: FoodGuard Architecture")

add_body_content(s3, [
    "FoodGuard is an enterprise-grade 4-class deep learning forensic system designed for high precision.",
    "Class Taxonomy:",
    "- 1. Authentic (Real): Genuine, unedited food photographs.",
    "- 2. Perfect AI: Pure Text-to-Image generations out of the box.",
    "- 3. Compressed AI: AI images degraded by JPEG compression to mimic social media sharing.",
    "- 4. Edited AI: Real images tampered with via localized AI inpainting (fraud objects)."
])

# ----------------- SLIDE 4: SYSTEM DESIGN -----------------
s4 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s4)
add_header_footer(s4, "System Design & Methodology")

# Using corporate boxes for flow
def flow_box(slide, x, y, w, h, title, subtitle):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLOR_BG
    shp.line.color.rgb = COLOR_ACCENT
    shp.line.width = Pt(2)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT
    p2.alignment = PP_ALIGN.CENTER

flow_box(s4, 1.0, 2.5, 3.0, 1.5, "1. Data Curation", "~191K Real Images\nKaggle: Food-101, etc.")
# Arrow (simple)
arr1 = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.0), Inches(0.6), Inches(0.4))
arr1.fill.solid()
arr1.fill.fore_color.rgb = COLOR_PRIMARY

flow_box(s4, 5.0, 2.5, 3.0, 1.5, "2. AI Synthesis", "RealVisXL / SDXL Inpaint\nGenerating Fraud Vectors")

arr2 = s4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.2), Inches(3.0), Inches(0.6), Inches(0.4))
arr2.fill.solid()
arr2.fill.fore_color.rgb = COLOR_PRIMARY

flow_box(s4, 9.0, 2.5, 3.0, 1.5, "3. Deep Training", "EfficientNet-B3 Backbone\nMixed Precision (AMP)")

arr3 = s4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.2), Inches(4.2), Inches(0.6), Inches(0.4))
arr3.fill.solid()
arr3.fill.fore_color.rgb = COLOR_PRIMARY

flow_box(s4, 9.0, 4.8, 3.0, 1.5, "4. Production UI", "Streamlit Dashboard\nError Level Analysis")


# ----------------- SLIDE 5: METRICS & EXCELLENCE -----------------
s5 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s5)
add_header_footer(s5, "Performance & Evaluation Metrics")

add_body_content(s5, [
    "Rigorous testing against a completely unseen holdout dataset (1,061 images).",
    "Overall Test Accuracy: 99.81%",
    "Target Benchmark: Zero Tolerance for False Positives on Real Food.",
    "- False Positive Rate (FPR) Achieved: 0.00%",
    "Model Specifications:",
    "- Framework: PyTorch & timm (EfficientNet-B3)",
    "- Optimization: AdamW with Cosine Annealing, weighted Cross-Entropy Loss.",
])

# Large emphasis box for metrics
box_metrics = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(4.5), Inches(4.0), Inches(1.5))
box_metrics.fill.solid()
box_metrics.fill.fore_color.rgb = COLOR_PRIMARY
tf_m = box_metrics.text_frame
p_m = tf_m.paragraphs[0]
p_m.text = "0.00% FPR\n99.8% Test Accuracy"
p_m.font.size = Pt(28)
p_m.font.bold = True
p_m.font.color.rgb = COLOR_WHITE
p_m.alignment = PP_ALIGN.CENTER

# ----------------- SLIDE 6: ELA FORENSICS -----------------
s6 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s6)
add_header_footer(s6, "Error Level Analysis (ELA)")

add_body_content(s6, [
    "Beyond Neural Networks: Interpretability via Image Forensics.",
    "How ELA Works:",
    "- Evaluates discrepancies in JPEG compression across the image grid.",
    "- Authentic images compress uniformly.",
    "- Localized tampering (e.g., edited insects) disrupts the compression grid.",
    "Strategic Advantage:",
    "- ELA serves as a deterministic cross-check alongside the probabilistic deep learning classification.",
    "- Instills enterprise trust by providing visual proof of manipulation."
])

# ----------------- SLIDE 7: FUTURE SCOPE -----------------
s7 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s7)
add_header_footer(s7, "Future Scope & Scale")

add_body_content(s7, [
    "1. Dual-Stream Frequency Architecture (RGB + FFT):",
    "- Extracting Fast Fourier Transform patterns to catch invisible AI structural artifacts.",
    "2. Explainability Engine (Grad-CAM):",
    "- Heatmaps directly identifying the pixel clusters triggering the AI flags.",
    "3. Enterprise Integration (API):",
    "- Packaging via FastAPI for seamless integration into delivery app backends (Zomato/Swiggy).",
    "4. Expanded Fraud Modalities:",
    "- Incorporating physical fraud (lighting mismatches, 3D structural inconsistencies)."
])


# ----------------- SLIDE 8: Q&A -----------------
s8 = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(s8)

# Center block
qt_rect = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5))
qt_rect.fill.solid()
qt_rect.fill.fore_color.rgb = COLOR_PRIMARY
qt_rect.line.fill.background()

tx_qa = s8.shapes.add_textbox(Inches(0), Inches(3.0), Inches(13.333), Inches(1.5))
tf_qa = tx_qa.text_frame
p_qa = tf_qa.paragraphs[0]
p_qa.text = "Thank You. Questions?"
p_qa.font.size = Pt(54)
p_qa.font.bold = True
p_qa.font.color.rgb = COLOR_WHITE
p_qa.alignment = PP_ALIGN.CENTER

# Team bottom
tx_end = s8.shapes.add_textbox(Inches(0), Inches(6.5), Inches(13.333), Inches(0.5))
tf_end = tx_end.text_frame
p_end = tf_end.paragraphs[0]
p_end.text = "FoodGuard | Raj | Aman | Rahul"
p_end.font.size = Pt(20)
p_end.font.color.rgb = COLOR_TEXT
p_end.alignment = PP_ALIGN.CENTER

prs.save('FoodGuard_MidSem_Presentation.pptx')
print("Corporate Presentation Generated successfully!")
