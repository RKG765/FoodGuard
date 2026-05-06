import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

def set_cell_border(cell, border_color="000000", border_width=Pt(1)):
    # Simple hack to set borders is not easily supported in base python-pptx, skipping complex borders
    pass

prs = Presentation()

# Constants
TITLE_COLOR = RGBColor(0, 51, 102)
TEXT_COLOR = RGBColor(50, 50, 50)
ACCENT_COLOR = RGBColor(0, 150, 200)
BG_COLOR = RGBColor(245, 247, 250)

def apply_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "FoodGuard"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
title.text_frame.paragraphs[0].font.size = Pt(54)
title.text_frame.paragraphs[0].font.bold = True

subtitle.text = "AI-Powered Food Fraud Detection System\nMid-Semester Presentation\n\nGroup Members:\nRaj Kumar\nAman Yadav\nRahul"
for p in subtitle.text_frame.paragraphs:
    p.font.color.rgb = TEXT_COLOR
    p.font.size = Pt(24)

# Function for Standard Slide
def create_slide(title_text, content_list, add_bg=True):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    if add_bg:
        apply_bg(slide)
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.color.rgb = TEXT_COLOR
        p.font.size = Pt(22)
        p.level = 0
    return slide

# Slide 2: Problem Statement
create_slide("Problem Statement", [
    "• Rise of AI-generated content makes visual verification difficult.",
    "• Common frauds involving food images:",
    "   - Claiming refunds using fake contaminants (cockroaches, mold).",
    "   - Generating fake high-quality restaurant reviews.",
    "• Binary classification (Real vs Fake) is insufficient and loses the 'why'.",
    "• Manual detection is nearly impossible due to realistic AI models (e.g., SDXL)."
])

# Slide 3: Our Solution - FoodGuard
create_slide("Our Solution: FoodGuard", [
    "• A 4-class deep learning forensic system replacing typical binary detectors.",
    "• Classes we identify:",
    "   1. Real: Genuine, unedited photographs.",
    "   2. Perfect AI: High-quality generated food (Text-to-Image).",
    "   3. Compressed AI: AI images degraded by JPEG compression & resizing.",
    "   4. Edited AI: Real images tampered via inpainting (e.g. insects inserted).",
    "• Goal: Achieve ≤ 5% False Positive Rate on authentic images."
])

# Slide 4: System Architecture (Diagram representation)
slide = prs.slides.add_slide(prs.slide_layouts[5])
apply_bg(slide)
title = slide.shapes.title
title.text = "System Architecture"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
title.text_frame.paragraphs[0].font.bold = True

# Create shapes for diagram instead of mermaid image
left = Inches(1)
top = Inches(2.5)
width = Inches(2)
height = Inches(1)

s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, width, height)
s1.text = "Data Collection\n(Kaggle Datasets)"
s1.fill.solid()
s1.fill.fore_color.rgb = ACCENT_COLOR

s2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), top, width, height)
s2.text = "AI Generation\n(RealVisXL + SDXL)"
s2.fill.solid()
s2.fill.fore_color.rgb = ACCENT_COLOR

s3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4.5), Inches(3), height)
s3.text = "4-Class Dataset Preparation"
s3.fill.solid()
s3.fill.fore_color.rgb = RGBColor(100, 150, 100)

s4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(4.5), Inches(3), height)
s4.text = "Training EfficientNet-B3"
s4.fill.solid()
s4.fill.fore_color.rgb = RGBColor(150, 100, 150)

s5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(6), Inches(3), height)
s5.text = "Streamlit UI & Inference"
s5.fill.solid()
s5.fill.fore_color.rgb = TITLE_COLOR

# Slide 5: Data & Generation Methodology
create_slide("Data & Generation Methodology", [
    "• Real Images (~191K): Sourced from Food-101, Indian Food, and UECFOOD256.",
    "• AI Generation Pipeline:",
    "   - RealVisXL V4.0 used for 'Perfect AI'.",
    "   - Compression scripts apply JPEG degradation for 'Compressed AI'.",
    "   - SDXL Inpainting applies fraud objects (insects, mold) for 'Edited AI'.",
    "• Final Train/Val/Test Split: 70% / 15% / 15%",
    "• Extensive validation avoids data leakage."
])

# Slide 6: Model & Training Details
create_slide("Model & Training Details", [
    "• Backbone: EfficientNet-B3 (combines accuracy and efficiency).",
    "• Input Resolution: 512x512 to preserve fine forensic artifacts.",
    "• Optimization:",
    "   - AdamW Optimizer with Cosine Annealing.",
    "   - Weighted Cross-Entropy Loss to heavily penalize False Positives on Real.",
    "   - Mixed Precision Training (AMP) used for faster training.",
    "• Calibration: Threshold sweeping executed post-training to ensure ≤ 5% FPR."
])

# Slide 7: Results & Evaluation
slide = prs.slides.add_slide(prs.slide_layouts[5])
apply_bg(slide)
title = slide.shapes.title
title.text = "Results & Evaluation"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
title.text_frame.paragraphs[0].font.bold = True

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "Test Accuracy: 99.81% | False Positive Rate: 0.00%"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TITLE_COLOR

# Add table for confusion matrix pseudo
rows = 5
cols = 4
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(2)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
headers = ['Class', 'Total', 'Correct', 'Accuracy']
for i, header in enumerate(headers):
    table.cell(0, i).text = header
    table.cell(0, i).text_frame.paragraphs[0].font.bold = True

data = [
    ['Real', '90', '90', '100%'],
    ['Perfect AI', '120', '119', '99.1%'],
    ['Compressed AI', '101', '101', '100%'],
    ['Edited AI', '750', '749', '99.8%']
]

for row_idx, row_data in enumerate(data):
    for col_idx, item in enumerate(row_data):
        table.cell(row_idx + 1, col_idx).text = item

# Slide 8: Error Level Analysis (ELA)
create_slide("Error Level Analysis (Forensics)", [
    "• Implemented Error Level Analysis (ELA) in the Streamlit Dashboard.",
    "• ELA highlights regions of an image with differing compression levels.",
    "• Perfect for detecting 'Edited AI' (inpainting):",
    "   - Inpainted objects (like insects) often introduce compression discrepancies.",
    "   - Bright spots in ELA reveal potential tampering.",
    "• ELA acts as an interpretable forensics tool alongside the neural network."
])

# Slide 9: Future Scope
create_slide("Future Scope", [
    "1. Dual-Stream Architecture:",
    "   - Integrating an FFT (Fast Fourier Transform) frequency analysis stream to better detect artifact patterns invisible in RGB.",
    "2. Explainability (Grad-CAM):",
    "   - Generating heatmaps so users can see exactly where the AI suspects tampering.",
    "3. Deployment & API API:",
    "   - Wrapping the model in FastAPI for integration with delivery apps (e.g., Zomato, Swiggy).",
    "4. Expanding the Fraud Database:",
    "   - Including physical tampering (like lighting mismatches) alongside text-to-image."
])

# Slide 10: Conclusion
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Thank You"
title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

subtitle.text = "Questions?\n\nFoodGuard Team:\nRaj Kumar | Aman यादव | Rahul"
subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

prs.save('FoodGuard_MidSem_Presentation.pptx')
print("Presentation generated successfully!")
